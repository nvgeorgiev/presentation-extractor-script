from pptx import Presentation
import glob

for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)